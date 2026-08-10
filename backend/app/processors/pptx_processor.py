import os

from pptx import Presentation

from app.processors.base_processor import BaseProcessor
from app.schemas.document_schema import DocumentSchema


class PPTXProcessor(BaseProcessor):
    """
    Processor responsible for extracting text and metadata
    from Microsoft PowerPoint (.pptx) presentations.
    """

    def open_document(self, file_path: str):
        """
        Open the PowerPoint presentation.
        """
        return Presentation(file_path)

    def extract_text(self, presentation) -> str:
        """
        Extract text from all slides.
        """

        slides_text = []

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text = shape.text.strip()

                    if text:

                        slides_text.append(text)

        return "\n".join(slides_text)

    def extract_metadata(self, presentation) -> dict:
        """
        Extract metadata from the PowerPoint presentation.
        """

        properties = presentation.core_properties

        return {

            "title": properties.title,

            "author": properties.author,

            "subject": properties.subject,

            "category": properties.category,

            "created": str(properties.created),

            "modified": str(properties.modified),

            "slides": str(len(presentation.slides))

        }

    def clean_text(self, text: str) -> str:
        """
        Clean extracted text by removing
        unnecessary whitespace and blank lines.
        """

        lines = []

        for line in text.splitlines():

            line = line.strip()

            if line:

                lines.append(line)

        return "\n".join(lines)

    def process(self, file_path: str) -> DocumentSchema:
        """
        Complete PPTX processing pipeline.

        Steps:
        1. Open presentation
        2. Extract text
        3. Extract metadata
        4. Clean text
        5. Generate document ID
        6. Return DocumentSchema
        """

        presentation = self.open_document(file_path)

        # Extract text
        text = self.extract_text(presentation)

        # Extract metadata
        metadata = self.extract_metadata(presentation)

        # Clean text
        cleaned_text = self.clean_text(text)

        # Generate unique document ID
        #
        # Example:
        # uploads/
        # 69759dba-d439-4cb2-8cd0-af880a4731b9.pptx
        #
        # becomes:
        # 69759dba-d439-4cb2-8cd0-af880a4731b9

        document_id = os.path.splitext(
            os.path.basename(file_path)
        )[0]

        # Create DocumentSchema
        result = DocumentSchema(

            document_id=document_id,

            filename=os.path.basename(file_path),

            file_type="pptx",

            text=cleaned_text,

            metadata=metadata

        )

        return result